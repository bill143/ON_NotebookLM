"""
Nexus Slide Engine — Feature 1C: PPTX/Slide Deck Generation
Codename: ESPERANTO

Provides:
- PowerPoint slide deck generation from structured content
- Branded slide templates (title, content, two-column, image)
- Auto-layout with heading hierarchy detection
- Export to PPTX format
"""

from __future__ import annotations

import io
import json
import re
from dataclasses import dataclass, field
from typing import Any, Optional

from loguru import logger

from src.infra.nexus_obs_tracing import traced
from src.exceptions import ValidationError


@dataclass
class SlideContent:
    """Content for a single slide."""
    layout: str = "content"      # "title", "content", "two_column", "section", "blank"
    title: str = ""
    body: str = ""
    bullets: list[str] = field(default_factory=list)
    left_content: str = ""       # For two-column layout
    right_content: str = ""
    notes: str = ""              # Speaker notes
    image_path: Optional[str] = None


@dataclass
class SlideConfig:
    """Slide deck configuration."""
    title: str = "Nexus Presentation"
    subtitle: str = ""
    author: str = "Nexus Notebook 11 LM"
    brand_color: str = "#6366f1"
    accent_color: str = "#8b5cf6"
    font_title: str = "Calibri"
    font_body: str = "Calibri"
    width_inches: float = 13.333
    height_inches: float = 7.5


@dataclass
class SlideResult:
    """Output from slide generation."""
    data: bytes
    filename: str
    slide_count: int
    file_size_bytes: int


class SlideEngine:
    """Generates PPTX slide decks from structured content."""

    @traced("slides.generate")
    async def generate(
        self,
        content: str | dict | list,
        config_dict: Optional[dict[str, Any]] = None,
    ) -> dict[str, Any]:
        """Generate a slide deck from content."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        except ImportError:
            raise ValidationError(
                "python-pptx is required for slide generation. Install with: pip install python-pptx"
            )

        config = SlideConfig(**(config_dict or {}))
        slides = self._parse_content(content)

        prs = Presentation()
        prs.slide_width = Inches(config.width_inches)
        prs.slide_height = Inches(config.height_inches)

        brand_rgb = self._hex_to_rgb(config.brand_color)
        accent_rgb = self._hex_to_rgb(config.accent_color)

        # Title slide
        title_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_layout)
        title_slide.shapes.title.text = config.title
        if title_slide.placeholders[1]:
            title_slide.placeholders[1].text = config.subtitle or config.author

        # Style title
        for paragraph in title_slide.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.color.rgb = RGBColor(*brand_rgb)

        # Content slides
        for slide_data in slides:
            if slide_data.layout == "section":
                layout = prs.slide_layouts[2]  # Section header
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_data.title
                for p in slide.shapes.title.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = RGBColor(*brand_rgb)
                        run.font.size = Pt(36)
                        run.font.bold = True
            else:
                layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_data.title

                # Style heading
                for p in slide.shapes.title.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = RGBColor(*brand_rgb)
                        run.font.size = Pt(28)
                        run.font.bold = True

                # Add body content
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()

                if slide_data.bullets:
                    for i, bullet in enumerate(slide_data.bullets):
                        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                        p.text = bullet
                        p.font.size = Pt(18)
                        p.space_after = Pt(8)
                        p.font.color.rgb = RGBColor(55, 65, 81)
                elif slide_data.body:
                    tf.paragraphs[0].text = slide_data.body
                    tf.paragraphs[0].font.size = Pt(18)
                    tf.paragraphs[0].font.color.rgb = RGBColor(55, 65, 81)

            # Speaker notes
            if slide_data.notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data.notes

        # Export
        buffer = io.BytesIO()
        prs.save(buffer)
        pptx_data = buffer.getvalue()

        safe_name = re.sub(r"[^\w\s-]", "", config.title.lower())
        safe_name = re.sub(r"[-\s]+", "_", safe_name)[:60]

        return {
            "type": "pptx",
            "data": pptx_data,
            "filename": f"{safe_name}.pptx",
            "slide_count": len(prs.slides),
            "file_size_bytes": len(pptx_data),
        }

    def _parse_content(self, content: str | dict | list) -> list[SlideContent]:
        """Parse various content formats into slide content."""
        slides: list[SlideContent] = []

        if isinstance(content, list):
            for item in content:
                if isinstance(item, dict):
                    slides.append(SlideContent(
                        layout=item.get("layout", "content"),
                        title=item.get("title", ""),
                        body=item.get("body", ""),
                        bullets=item.get("bullets", []),
                        notes=item.get("notes", ""),
                    ))
            return slides

        if isinstance(content, dict):
            content = json.dumps(content, indent=2)

        # Parse markdown-like content
        text = str(content)
        sections = re.split(r"(?=^#{1,2}\s+)", text, flags=re.MULTILINE)

        for section in sections:
            section = section.strip()
            if not section:
                continue

            # Extract heading
            heading_match = re.match(r"^(#{1,3})\s+(.+?)$", section, re.MULTILINE)
            if heading_match:
                level = len(heading_match.group(1))
                title = heading_match.group(2).strip()
                body_text = section[heading_match.end():].strip()

                if level == 1:
                    slides.append(SlideContent(layout="section", title=title))

                # Extract bullets
                bullets = []
                remaining_text = []
                for line in body_text.split("\n"):
                    line = line.strip()
                    if re.match(r"^[-*•]\s+", line):
                        bullets.append(re.sub(r"^[-*•]\s+", "", line))
                    elif line:
                        remaining_text.append(line)

                if bullets or remaining_text:
                    slides.append(SlideContent(
                        title=title,
                        body="\n".join(remaining_text) if not bullets else "",
                        bullets=bullets,
                    ))
            else:
                # No heading — treat as content slide
                lines = [line.strip() for line in section.split("\n") if line.strip()]
                if lines:
                    slides.append(SlideContent(
                        title=lines[0][:80],
                        body="\n".join(lines[1:]),
                    ))

        return slides or [SlideContent(title="Content", body=str(content)[:500])]

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
        h = hex_color.lstrip("#")
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


# Global singleton
slide_engine = SlideEngine()
