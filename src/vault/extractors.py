"""
Nexus Vault — File Content Extractors
Codename: ESPERANTO

Extracts readable content from uploaded files for AI classification.
Each extractor returns a dict with 'content' (str) and 'file_metadata' (dict).
Graceful degradation: extraction failures return filename-only metadata, never crash.
"""

from __future__ import annotations

import csv
import io
import zipfile
from pathlib import Path
from typing import Any

from loguru import logger


# ── Extraction Result ───────────────────────────────────────


def _make_result(
    content: str,
    file_metadata: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Build a standardized extraction result."""
    return {
        "content": content[:50_000],  # Cap at 50k chars for AI context
        "file_metadata": file_metadata or {},
    }


# ── PDF Extractor ───────────────────────────────────────────


def extract_pdf(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Extract text from PDF using pypdf."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(file_bytes))
        pages_text = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages_text.append(f"[Page {i + 1}]\n{text}")

        content = "\n\n".join(pages_text)
        return _make_result(
            content=content or f"[PDF file: {filename} — no extractable text]",
            file_metadata={
                "page_count": len(reader.pages),
                "has_text": bool(content.strip()),
                "file_size_bytes": len(file_bytes),
            },
        )
    except Exception as e:
        logger.warning(f"PDF extraction failed for {filename}: {e}")
        return _make_result(
            content=f"[PDF file: {filename} — extraction failed]",
            file_metadata={"extraction_error": str(e)[:200]},
        )


# ── DWG/DXF Extractor ──────────────────────────────────────


def extract_dwg(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Extract metadata from DWG/DXF CAD files."""
    ext = Path(filename).suffix.lower()
    metadata: dict[str, Any] = {
        "file_type": ext.lstrip("."),
        "file_size_bytes": len(file_bytes),
    }

    # Attempt basic DXF text extraction (DXF is ASCII-based)
    if ext == ".dxf":
        try:
            text = file_bytes.decode("utf-8", errors="replace")
            # Extract layer names from DXF
            layers: list[str] = []
            lines = text.split("\n")
            for i, line in enumerate(lines):
                if line.strip() == "LAYER" and i + 2 < len(lines):
                    name_line = lines[i + 2].strip()
                    if name_line and name_line not in layers:
                        layers.append(name_line)
            metadata["layers"] = layers[:50]
            metadata["layer_count"] = len(layers)

            # Try to find title block text
            title_block_hints = []
            for line in lines:
                stripped = line.strip()
                if any(kw in stripped.upper() for kw in ["PROJECT", "DRAWING", "SHEET", "REV"]):
                    if 3 < len(stripped) < 200:
                        title_block_hints.append(stripped)
            metadata["title_block_hints"] = title_block_hints[:10]

            content_parts = [f"CAD Drawing: {filename}"]
            if layers:
                content_parts.append(f"Layers ({len(layers)}): {', '.join(layers[:20])}")
            if title_block_hints:
                content_parts.append(f"Title block info: {'; '.join(title_block_hints[:5])}")

            return _make_result(
                content="\n".join(content_parts),
                file_metadata=metadata,
            )
        except Exception as e:
            logger.debug(f"DXF parsing failed for {filename}: {e}")

    # DWG binary — extract what we can from filename
    return _make_result(
        content=(
            f"CAD Drawing file: {filename}\n"
            f"File type: {ext}\n"
            f"File size: {len(file_bytes):,} bytes\n"
            f"Note: Binary CAD format — classify based on filename and project context."
        ),
        file_metadata=metadata,
    )


# ── Spreadsheet Extractor ──────────────────────────────────


def extract_spreadsheet(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Extract content from XLS, XLSX, and CSV files."""
    ext = Path(filename).suffix.lower()
    metadata: dict[str, Any] = {
        "file_type": ext.lstrip("."),
        "file_size_bytes": len(file_bytes),
    }

    # CSV
    if ext == ".csv":
        return _extract_csv(file_bytes, filename, metadata)

    # XLSX / XLS via openpyxl
    if ext in (".xlsx", ".xls"):
        return _extract_xlsx(file_bytes, filename, metadata)

    return _make_result(
        content=f"Spreadsheet file: {filename}",
        file_metadata=metadata,
    )


def _extract_csv(
    file_bytes: bytes, filename: str, metadata: dict[str, Any]
) -> dict[str, Any]:
    """Extract CSV content."""
    try:
        text = file_bytes.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)

        if not rows:
            return _make_result(
                content=f"Empty CSV file: {filename}",
                file_metadata=metadata,
            )

        headers = rows[0] if rows else []
        metadata["column_headers"] = headers
        metadata["row_count"] = len(rows) - 1  # Exclude header
        metadata["column_count"] = len(headers)

        content_parts = [
            f"CSV File: {filename}",
            f"Columns ({len(headers)}): {', '.join(headers[:20])}",
            f"Row count: {len(rows) - 1}",
            "",
            "First 10 rows:",
        ]
        for row in rows[1:11]:
            content_parts.append(" | ".join(str(cell)[:50] for cell in row))

        return _make_result(
            content="\n".join(content_parts),
            file_metadata=metadata,
        )
    except Exception as e:
        logger.warning(f"CSV extraction failed for {filename}: {e}")
        return _make_result(
            content=f"CSV file: {filename} — extraction failed",
            file_metadata={**metadata, "extraction_error": str(e)[:200]},
        )


def _extract_xlsx(
    file_bytes: bytes, filename: str, metadata: dict[str, Any]
) -> dict[str, Any]:
    """Extract XLSX content using openpyxl."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        sheet_names = wb.sheetnames
        metadata["sheet_names"] = sheet_names
        metadata["sheet_count"] = len(sheet_names)

        content_parts = [
            f"Excel File: {filename}",
            f"Sheets ({len(sheet_names)}): {', '.join(sheet_names[:10])}",
        ]

        for sheet_name in sheet_names[:3]:  # First 3 sheets
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(max_row=11, values_only=True))
            if not rows:
                continue

            headers = [str(cell or "") for cell in rows[0]]
            row_count = ws.max_row or 0
            metadata[f"sheet_{sheet_name}_headers"] = headers
            metadata[f"sheet_{sheet_name}_rows"] = row_count

            content_parts.append(f"\n--- Sheet: {sheet_name} ---")
            content_parts.append(f"Columns: {', '.join(headers[:15])}")
            content_parts.append(f"Rows: ~{row_count}")
            content_parts.append("First 10 rows:")
            for row in rows[1:11]:
                cells = [str(cell or "")[:50] for cell in row]
                content_parts.append(" | ".join(cells))

        wb.close()
        return _make_result(
            content="\n".join(content_parts),
            file_metadata=metadata,
        )
    except Exception as e:
        logger.warning(f"XLSX extraction failed for {filename}: {e}")
        return _make_result(
            content=f"Excel file: {filename} — extraction failed",
            file_metadata={**metadata, "extraction_error": str(e)[:200]},
        )


# ── XER Extractor (Primavera P6) ───────────────────────────


def extract_xer(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Extract metadata from Primavera P6 XER schedule files."""
    metadata: dict[str, Any] = {
        "file_type": "xer",
        "file_size_bytes": len(file_bytes),
    }

    try:
        text = file_bytes.decode("utf-8", errors="replace")
        lines = text.split("\n")

        project_name = None
        data_date = None
        activity_count = 0

        current_table = None
        for line in lines:
            stripped = line.strip()

            # Table headers start with %T
            if stripped.startswith("%T"):
                current_table = stripped[2:].strip()
                continue

            # Data rows start with %R
            if stripped.startswith("%R") and current_table:
                fields = stripped[2:].split("\t")

                if current_table == "PROJECT" and len(fields) > 1:
                    project_name = fields[1] if len(fields) > 1 else None

                if current_table == "PROJWBS" and not project_name and len(fields) > 1:
                    project_name = fields[1]

                if current_table == "TASK":
                    activity_count += 1

                if current_table == "PROJECT" and len(fields) > 3:
                    data_date = fields[3] if len(fields) > 3 else None

        metadata["project_name"] = project_name
        metadata["data_date"] = data_date
        metadata["activity_count"] = activity_count

        content = (
            f"Primavera P6 Schedule (XER): {filename}\n"
            f"Project: {project_name or 'Unknown'}\n"
            f"Data Date: {data_date or 'Unknown'}\n"
            f"Activity Count: {activity_count}\n"
            f"This is a construction CPM schedule export from Oracle Primavera P6."
        )

        return _make_result(content=content, file_metadata=metadata)
    except Exception as e:
        logger.warning(f"XER extraction failed for {filename}: {e}")
        return _make_result(
            content=f"Primavera P6 schedule file: {filename}",
            file_metadata={**metadata, "extraction_error": str(e)[:200]},
        )


# ── RVT Extractor (Revit) ──────────────────────────────────


def extract_rvt(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Extract metadata from Revit (RVT) files."""
    # RVT is a binary OLE compound document — extract what we can
    name_lower = filename.lower()

    discipline = "unknown"
    discipline_keywords = {
        "architectural": ["arch", "a-", "a_"],
        "structural": ["struct", "s-", "s_"],
        "mechanical": ["mech", "m-", "m_", "hvac"],
        "electrical": ["elec", "e-", "e_"],
        "plumbing": ["plumb", "p-", "p_"],
        "civil": ["civil", "c-", "c_", "site"],
        "landscape": ["land", "l-", "l_"],
    }
    for disc, keywords in discipline_keywords.items():
        if any(kw in name_lower for kw in keywords):
            discipline = disc
            break

    return _make_result(
        content=(
            f"Revit BIM Model: {filename}\n"
            f"Discipline (from filename): {discipline}\n"
            f"File size: {len(file_bytes):,} bytes\n"
            f"This is an Autodesk Revit Building Information Model file."
        ),
        file_metadata={
            "file_type": "rvt",
            "discipline": discipline,
            "file_size_bytes": len(file_bytes),
        },
    )


# ── Image Extractor ─────────────────────────────────────────


def extract_image(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Build description for image files (PNG, JPG, etc.) for vision API."""
    ext = Path(filename).suffix.lower()
    return _make_result(
        content=(
            f"Image file: {filename}\n"
            f"Format: {ext.lstrip('.')}\n"
            f"File size: {len(file_bytes):,} bytes\n"
            f"[This image should be analyzed by the vision API for content description]"
        ),
        file_metadata={
            "file_type": ext.lstrip("."),
            "file_size_bytes": len(file_bytes),
            "requires_vision": True,
        },
    )


# ── ZIP Extractor ───────────────────────────────────────────


def extract_zip(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """List contents of a ZIP archive for classification."""
    metadata: dict[str, Any] = {
        "file_type": "zip",
        "file_size_bytes": len(file_bytes),
    }

    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            file_list = zf.namelist()
            metadata["file_count"] = len(file_list)
            metadata["contained_files"] = file_list[:100]

            # Analyze file types in the archive
            extensions: dict[str, int] = {}
            for f in file_list:
                ext = Path(f).suffix.lower()
                if ext:
                    extensions[ext] = extensions.get(ext, 0) + 1
            metadata["extension_breakdown"] = extensions

            content_parts = [
                f"ZIP Archive: {filename}",
                f"Files: {len(file_list)}",
                f"Extensions: {dict(sorted(extensions.items(), key=lambda x: -x[1]))}",
                "",
                "Contents:",
            ]
            for f in file_list[:50]:
                content_parts.append(f"  {f}")
            if len(file_list) > 50:
                content_parts.append(f"  ... and {len(file_list) - 50} more files")

            return _make_result(
                content="\n".join(content_parts),
                file_metadata=metadata,
            )
    except Exception as e:
        logger.warning(f"ZIP extraction failed for {filename}: {e}")
        return _make_result(
            content=f"ZIP archive: {filename} — extraction failed",
            file_metadata={**metadata, "extraction_error": str(e)[:200]},
        )


# ── Media Extractor (MP4, M4A) ─────────────────────────────


def extract_media(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Extract metadata from media files."""
    ext = Path(filename).suffix.lower()
    return _make_result(
        content=(
            f"Media file: {filename}\n"
            f"Format: {ext.lstrip('.')}\n"
            f"File size: {len(file_bytes):,} bytes\n"
            f"Classify based on filename and project context."
        ),
        file_metadata={
            "file_type": ext.lstrip("."),
            "file_size_bytes": len(file_bytes),
        },
    )


# ── PPTX Extractor ─────────────────────────────────────────


def extract_pptx(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Extract slide content from PowerPoint files."""
    metadata: dict[str, Any] = {
        "file_type": "pptx",
        "file_size_bytes": len(file_bytes),
    }

    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        slide_count = len(prs.slides)
        metadata["slide_count"] = slide_count

        content_parts = [
            f"PowerPoint Presentation: {filename}",
            f"Slides: {slide_count}",
        ]

        for i, slide in enumerate(prs.slides):
            if i >= 5:  # First 5 slides
                break
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                content_parts.append(f"\n--- Slide {i + 1} ---")
                content_parts.append("\n".join(texts[:10]))

        return _make_result(
            content="\n".join(content_parts),
            file_metadata=metadata,
        )
    except Exception as e:
        logger.warning(f"PPTX extraction failed for {filename}: {e}")
        return _make_result(
            content=f"PowerPoint file: {filename} — extraction failed",
            file_metadata={**metadata, "extraction_error": str(e)[:200]},
        )


# ── DOCX Extractor ─────────────────────────────────────────


def extract_docx(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Extract text from Word documents."""
    metadata: dict[str, Any] = {
        "file_type": "docx",
        "file_size_bytes": len(file_bytes),
    }

    try:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        metadata["paragraph_count"] = len(paragraphs)

        content = f"Word Document: {filename}\n\n" + "\n".join(paragraphs)
        return _make_result(content=content, file_metadata=metadata)
    except Exception as e:
        logger.warning(f"DOCX extraction failed for {filename}: {e}")
        return _make_result(
            content=f"Word document: {filename} — extraction failed",
            file_metadata={**metadata, "extraction_error": str(e)[:200]},
        )


# ── GeoJSON Extractor ──────────────────────────────────────


def extract_geojson(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Extract metadata from GeoJSON files."""
    import json

    metadata: dict[str, Any] = {
        "file_type": "geojson",
        "file_size_bytes": len(file_bytes),
    }

    try:
        data = json.loads(file_bytes.decode("utf-8"))
        features = data.get("features", [])
        metadata["feature_count"] = len(features)
        metadata["geometry_types"] = list(
            {f.get("geometry", {}).get("type", "unknown") for f in features[:100]}
        )

        content = (
            f"GeoJSON File: {filename}\n"
            f"Features: {len(features)}\n"
            f"Geometry types: {', '.join(metadata['geometry_types'])}\n"
            f"This is a geospatial data file — likely a survey or site plan."
        )
        return _make_result(content=content, file_metadata=metadata)
    except Exception as e:
        logger.warning(f"GeoJSON extraction failed for {filename}: {e}")
        return _make_result(
            content=f"GeoJSON file: {filename}",
            file_metadata={**metadata, "extraction_error": str(e)[:200]},
        )


# ── Fallback Extractor ─────────────────────────────────────


def extract_fallback(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """Fallback extractor — filename, extension, and file size metadata only."""
    ext = Path(filename).suffix.lower()
    return _make_result(
        content=(
            f"File: {filename}\n"
            f"Extension: {ext or 'none'}\n"
            f"File size: {len(file_bytes):,} bytes\n"
            f"No specialized extractor available — classify based on filename and context."
        ),
        file_metadata={
            "file_type": ext.lstrip(".") if ext else "unknown",
            "file_size_bytes": len(file_bytes),
        },
    )


# ── Extractor Router ───────────────────────────────────────


_EXTENSION_MAP: dict[str, Any] = {
    ".pdf": extract_pdf,
    ".dwg": extract_dwg,
    ".dxf": extract_dwg,
    ".xls": extract_spreadsheet,
    ".xlsx": extract_spreadsheet,
    ".csv": extract_spreadsheet,
    ".xer": extract_xer,
    ".rvt": extract_rvt,
    ".rfa": extract_rvt,
    ".ifc": extract_rvt,
    ".png": extract_image,
    ".jpg": extract_image,
    ".jpeg": extract_image,
    ".tiff": extract_image,
    ".tif": extract_image,
    ".bmp": extract_image,
    ".webp": extract_image,
    ".zip": extract_zip,
    ".mp4": extract_media,
    ".m4a": extract_media,
    ".mp3": extract_media,
    ".wav": extract_media,
    ".pptx": extract_pptx,
    ".ppt": extract_pptx,
    ".docx": extract_docx,
    ".doc": extract_docx,
    ".geojson": extract_geojson,
}


def extract_content(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """
    Route file to the appropriate extractor based on extension.

    Returns:
        {"content": str, "file_metadata": dict}
    """
    ext = Path(filename).suffix.lower()
    extractor = _EXTENSION_MAP.get(ext, extract_fallback)

    logger.debug(f"Extracting content from {filename} using {extractor.__name__}")

    try:
        return extractor(file_bytes, filename)
    except Exception as e:
        logger.error(f"Extraction failed for {filename}: {e}")
        return extract_fallback(file_bytes, filename)
